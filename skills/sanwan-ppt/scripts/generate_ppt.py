#!/usr/bin/env python3
"""
三万同款板书风格 PPT 生成脚本
用法：python3 generate_ppt.py '<JSON格式slides数组>' [输出文件名]

slides JSON 格式：
[
  {"title": "页面标题", "prompt_body": "页面内容描述（不含STYLE前缀）", "mascot_position": "bottom-left", "mascot_mood": "excited"}
]
"""

import sys
import json
import base64
import os
import requests

# ── 读取 API Key ──────────────────────────────────────────────
def get_api_key():
    config_path = os.path.expanduser("~/.openclaw/openclaw.json")
    with open(config_path) as f:
        cfg = json.load(f)
    return cfg["models"]["providers"]["deepv-easyclaw"]["apiKey"]

# ── STYLE 前缀 ────────────────────────────────────────────────
STYLE = """Whiteboard filling the ENTIRE 16:9 frame, all four silver/gray magnetic \
frame borders fully visible at image edges, zero background, zero office or room visible. \
Clean white whiteboard surface. All text rendered in elegant Chinese hard-pen fountain pen \
handwriting calligraphy (硬笔书法) — precise clean strokes, ink variation, flowing and \
neat, fine pen line quality, NOT thick marker, NOT digital font. Illustrations are \
hand-drawn cartoon style: bold black marker outlines (like Copic multiliner), filled with \
vivid colored markers (Copic/Prismacolor style), layered shadows and highlights for depth, \
NOT flat vector, NOT watercolor wash — solid marker color with visible stroke direction. \
MANDATORY MASCOT on every single page without exception: a super cute chibi Labrador \
retriever wearing a red lobster-claw hat, big sparkling eyes, rosy cheeks, fluffy golden \
fur, hand-drawn marker style — expression matching the page mood (curious, excited, \
thinking, celebrating, etc.). DO NOT omit the Labrador mascot under any circumstances. \
Annotation marks in red or black hand-drawn pen (arrows, underlines, ✕, ✓). 16:9 widescreen."""

def build_prompt(prompt_body: str) -> str:
    return STYLE + "\n\n" + prompt_body

# ── 生成单张图片 ──────────────────────────────────────────────
def generate_image(api_key: str, prompt: str, output_path: str):
    resp = requests.post(
        "https://api.easyclaw.work/api/v1/images/generate",
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        },
        json={
            "prompt": prompt,
            "aspect_ratio": "16:9",
            "resolution": "2K"
        },
        timeout=120
    ).json()

    if "image_url" not in resp:
        raise RuntimeError(f"图片生成失败: {json.dumps(resp, ensure_ascii=False)}")

    b64 = resp["image_url"].split(",", 1)[1]
    with open(output_path, "wb") as f:
        f.write(base64.b64decode(b64))
    print(f"  [OK] 图片已保存: {output_path}", flush=True)

# ── 组装 PPTX ────────────────────────────────────────────────
def build_pptx(image_paths: list, output_path: str):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for img_path in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
        slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

    prs.save(output_path)
    print(f"\n[OK] PPT 已保存: {output_path}", flush=True)

# ── 主流程 ───────────────────────────────────────────────────
def main():
    if len(sys.argv) < 2:
        print("用法: python3 generate_ppt.py '<slides_json>' [output.pptx]")
        sys.exit(1)

    slides = json.loads(sys.argv[1])
    output_path = sys.argv[2] if len(sys.argv) > 2 else "/tmp/sanwan_output.pptx"

    api_key = get_api_key()
    image_paths = []

    print(f"共 {len(slides)} 页，开始生成...\n", flush=True)

    for i, slide in enumerate(slides, 1):
        print(f"[{i}/{len(slides)}] 生成: {slide.get('title', f'第{i}页')}", flush=True)
        prompt = build_prompt(slide["prompt_body"])
        img_path = f"/tmp/sanwan_slide_{i:02d}.png"
        generate_image(api_key, prompt, img_path)
        image_paths.append(img_path)

    build_pptx(image_paths, output_path)
    # 输出路径供调用方捕获
    print(f"OUTPUT_PATH={output_path}", flush=True)

if __name__ == "__main__":
    main()
